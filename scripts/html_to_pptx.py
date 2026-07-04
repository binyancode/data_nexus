"""
Convert semantic_tool_layer_flow_ppt.html → PPT with embedded SVGs + text boxes.

Usage:
    Activate root .venv, then:
    python tools/html_to_pptx.py

Output:
    doc/design/semantic_tool_layer_flow_ppt.pptx
    doc/design/svg/slide_XX.svg  (extracted SVGs)
"""

import os
import re
import io
from pathlib import Path
from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree

# --- Paths ---
BASE_DIR = Path(__file__).resolve().parent.parent
HTML_PATH = BASE_DIR / "doc" / "design" / "semantic_tool_layer_flow_ppt.html"
SVG_DIR = BASE_DIR / "doc" / "design" / "svg"
OUTPUT_PATH = BASE_DIR / "doc" / "design" / "semantic_tool_layer_flow_ppt.pptx"

# --- Colors (matching HTML theme) ---
BG_COLOR = RGBColor(0x14, 0x1E, 0x30)
SURFACE_COLOR = RGBColor(0x1C, 0x2A, 0x42)
TEXT_COLOR = RGBColor(0xEA, 0xF0, 0xF6)
TEXT_MUTED = RGBColor(0xB0, 0xBE, 0xC5)
ACCENT_BLUE = RGBColor(0x4F, 0xAC, 0xFE)
ACCENT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE)

ICON_COLORS = {
    "blue": ACCENT_BLUE,
    "orange": ACCENT_ORANGE,
    "purple": ACCENT_PURPLE,
    "green": ACCENT_GREEN,
    "cyan": ACCENT_CYAN,
}

# Slide dimensions: 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def extract_slides(html_path: Path) -> list[dict]:
    """Parse the HTML and extract slide info: title, subtitle, svg content, notes."""
    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    # Extract the <style> block for inlining into SVGs
    style_tag = soup.find("style")
    style_text = style_tag.string if style_tag else ""

    slides_data = []
    slide_divs = soup.find_all("div", class_="slide")

    for i, slide_div in enumerate(slide_divs):
        slide_info = {
            "index": i + 1,
            "title": "",
            "subtitle": "",
            "svg": None,
            "notes": [],
            "is_title_slide": "title-slide" in slide_div.get("class", []),
            "has_swimlane": False,
            "has_compare": False,
        }

        # Title
        h2 = slide_div.find("h2")
        if h2:
            slide_info["title"] = h2.get_text()

        h1 = slide_div.find("h1")
        if h1:
            slide_info["title"] = h1.get_text()

        # Subtitle
        sub = slide_div.find("div", class_="phase-subtitle")
        if sub:
            slide_info["subtitle"] = sub.get_text()

        h2_tag = slide_div.find("h2")
        if slide_info["is_title_slide"]:
            h2_sub = slide_div.find("h2")
            if h2_sub:
                slide_info["subtitle"] = h2_sub.get_text()

        # SVG
        svg = slide_div.find("svg")
        if svg:
            # Add xmlns if missing
            if not svg.get("xmlns"):
                svg["xmlns"] = "http://www.w3.org/2000/svg"

            # Inject a minimal style for font-family into SVG
            svg_style = soup.new_tag("style")
            svg_style.string = """
                text { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; }
                text[fill="#7e8fa6"] { fill: #b0bec5; }
            """
            svg.insert(0, svg_style)
            slide_info["svg"] = str(svg)

        # Swimlane/compare detection
        if slide_div.find("div", class_="swimlane-flow"):
            slide_info["has_swimlane"] = True
        if slide_div.find("div", class_="compare-grid"):
            slide_info["has_compare"] = True

        # Notes
        notes_div = slide_div.find("div", class_="slide-notes")
        if notes_div:
            for h4 in notes_div.find_all("h4"):
                note_entry = {"heading": h4.get_text(), "color": "blue", "paragraphs": []}

                # Detect color from icon
                icon = h4.find("span", class_="note-icon")
                if icon:
                    classes = icon.get("class", [])
                    for c in classes:
                        if c in ICON_COLORS:
                            note_entry["color"] = c
                            break

                # Collect siblings until next h4 or note-divider
                for sibling in h4.next_siblings:
                    if isinstance(sibling, Tag):
                        if sibling.name == "h4":
                            break
                        if "note-divider" in sibling.get("class", []):
                            break
                        if sibling.name == "p":
                            note_entry["paragraphs"].append(sibling.get_text())
                        if sibling.name == "ul":
                            for li in sibling.find_all("li"):
                                note_entry["paragraphs"].append("• " + li.get_text())

                slide_info["notes"].append(note_entry)

        slides_data.append(slide_info)

    return slides_data


def save_svgs(slides_data: list[dict], svg_dir: Path):
    """Save extracted SVGs as standalone .svg files."""
    svg_dir.mkdir(parents=True, exist_ok=True)
    for slide in slides_data:
        if slide["svg"]:
            svg_path = svg_dir / f"slide_{slide['index']:02d}.svg"
            # Wrap in proper XML declaration
            svg_content = '<?xml version="1.0" encoding="UTF-8"?>\n' + slide["svg"]
            with open(svg_path, "w", encoding="utf-8") as f:
                f.write(svg_content)
            print(f"  Saved: {svg_path.name}")


def add_svg_picture(slide, svg_path: str, left, top, width, height):
    """Embed SVG directly into PowerPoint using low-level OPC manipulation.
    PowerPoint 2019+ / Office 365 supports SVG natively."""
    from pptx.parts.image import ImagePart
    from pptx.package import Package

    # Read SVG bytes
    with open(svg_path, "rb") as f:
        svg_bytes = f.read()

    # Create the image part manually
    slide_part = slide.part
    partname = slide_part.package.next_image_partname(".svg")

    # Add SVG as a binary part with correct content type
    from pptx.opc.package import Part

    image_part = Part(partname, "image/svg+xml", slide_part.package, blob=svg_bytes)
    rId = slide_part.relate_to(image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")

    # Create pic element
    pic_id = len(slide.shapes) + 1
    pic_xml = (
        f'<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <p:nvPicPr>'
        f'    <p:cNvPr id="{pic_id}" name="SVG {pic_id}"/>'
        f'    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'    <p:nvPr/>'
        f'  </p:nvPicPr>'
        f'  <p:blipFill>'
        f'    <a:blip r:embed="{rId}"/>'
        f'    <a:stretch><a:fillRect/></a:stretch>'
        f'  </p:blipFill>'
        f'  <p:spPr>'
        f'    <a:xfrm>'
        f'      <a:off x="{int(left)}" y="{int(top)}"/>'
        f'      <a:ext cx="{int(width)}" cy="{int(height)}"/>'
        f'    </a:xfrm>'
        f'    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'  </p:spPr>'
        f'</p:pic>'
    )
    pic_element = parse_xml(pic_xml)
    slide.shapes._spTree.append(pic_element)


def set_slide_bg(slide, color: RGBColor):
    """Set solid fill background for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_notes_textbox(slide, notes: list[dict], left, top, width, max_height):
    """Add the notes panel as a single formatted text box."""
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(left, top, width, max_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Pt(16)
    tf.margin_bottom = Pt(16)
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(18)

    # Add background fill to the shape
    sp = txBox._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        from pptx.oxml import parse_xml
        spPr = parse_xml(f'<p:spPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        sp.append(spPr)

    # Use shape fill
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = SURFACE_COLOR

    # Border
    line = txBox.line
    line.color.rgb = RGBColor(0x2E, 0x40, 0x60)
    line.width = Pt(1)

    first_para = True
    for note in notes:
        # Heading
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            # Add spacing line
            p = tf.add_paragraph()
            p.space_before = Pt(8)

            p = tf.add_paragraph()

        run = p.add_run()
        run.text = note["heading"]
        run.font.size = Pt(12)
        run.font.bold = True
        color_key = note.get("color", "blue")
        run.font.color.rgb = ICON_COLORS.get(color_key, ACCENT_BLUE)

        # Body paragraphs
        for para_text in note["paragraphs"]:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = para_text
            run.font.size = Pt(10)
            run.font.color.rgb = TEXT_MUTED
            p.space_before = Pt(2)

    return txBox


def build_pptx(slides_data: list[dict], svg_dir: Path, output_path: Path):
    """Generate the PPT file."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, BG_COLOR)

        if slide_data["is_title_slide"]:
            # Title slide
            add_text_box(slide,
                         left=Inches(1), top=Inches(2.5),
                         width=Inches(11.3), height=Inches(1.5),
                         text=slide_data["title"],
                         font_size=48, color=ACCENT_BLUE, bold=True,
                         alignment=PP_ALIGN.CENTER)
            if slide_data["subtitle"]:
                add_text_box(slide,
                             left=Inches(1), top=Inches(4.0),
                             width=Inches(11.3), height=Inches(1),
                             text=slide_data["subtitle"],
                             font_size=18, color=TEXT_MUTED,
                             alignment=PP_ALIGN.CENTER)
        else:
            # Header: title
            add_text_box(slide,
                         left=Inches(0.8), top=Inches(0.3),
                         width=Inches(10), height=Inches(0.7),
                         text=slide_data["title"],
                         font_size=24, color=TEXT_COLOR, bold=True)

            # Subtitle
            if slide_data["subtitle"]:
                add_text_box(slide,
                             left=Inches(0.8), top=Inches(0.9),
                             width=Inches(10), height=Inches(0.5),
                             text=slide_data["subtitle"],
                             font_size=12, color=TEXT_MUTED)

            # SVG image (embedded directly for PowerPoint 2019+/365)
            svg_file = svg_dir / f"slide_{slide_data['index']:02d}.svg"
            if svg_file.exists():
                if slide_data["notes"]:
                    # Layout: SVG left + notes right
                    add_svg_picture(slide, str(svg_file),
                                    left=Inches(0.3), top=Inches(1.5),
                                    width=Inches(9.5), height=Inches(5.7))

                    # Notes panel on right
                    notes_left = Inches(9.9)
                    notes_top = Inches(1.5)
                    notes_width = Inches(3.2)
                    notes_height = Inches(5.7)

                    add_notes_textbox(slide, slide_data["notes"],
                                      notes_left, notes_top, notes_width, notes_height)
                else:
                    # Full width SVG
                    add_svg_picture(slide, str(svg_file),
                                    left=Inches(0.3), top=Inches(1.5),
                                    width=Inches(12.7), height=Inches(5.7))
            elif slide_data["has_swimlane"] or slide_data["has_compare"]:
                # For swimlane/compare slides without SVG, add a placeholder note
                add_text_box(slide,
                             left=Inches(0.8), top=Inches(2),
                             width=Inches(11), height=Inches(1),
                             text="[See HTML version for full swimlane/comparison layout]",
                             font_size=14, color=TEXT_MUTED,
                             alignment=PP_ALIGN.CENTER)

        # Slide number
        add_text_box(slide,
                     left=Inches(11.5), top=Inches(7.0),
                     width=Inches(1.5), height=Inches(0.4),
                     text=f"{slide_data['index']:02d} / {len(slides_data):02d}",
                     font_size=10, color=TEXT_MUTED,
                     alignment=PP_ALIGN.RIGHT)

        # Brand corner
        add_text_box(slide,
                     left=Inches(10.5), top=Inches(0.15),
                     width=Inches(2.5), height=Inches(0.4),
                     text="BEONE AIBI",
                     font_size=9, color=TEXT_MUTED,
                     alignment=PP_ALIGN.RIGHT)

    prs.save(str(output_path))
    print(f"\n✓ PPT saved to: {output_path}")


def main():
    print("=" * 50)
    print("HTML → PPT Converter (SVG embedded)")
    print("=" * 50)

    print(f"\nParsing: {HTML_PATH.name}")
    slides_data = extract_slides(HTML_PATH)
    print(f"  Found {len(slides_data)} slides")

    print(f"\nExtracting SVGs to: {SVG_DIR}/")
    save_svgs(slides_data, SVG_DIR)

    print(f"\nBuilding PPT...")
    build_pptx(slides_data, SVG_DIR, OUTPUT_PATH)

    print(f"\nDone! Files created:")
    print(f"  • {OUTPUT_PATH}")
    print(f"  • {SVG_DIR}/ ({sum(1 for s in slides_data if s['svg'])} SVG files)")
    print(f"\nTip: In PowerPoint, right-click SVG → 'Convert to Shape' for full editability")


if __name__ == "__main__":
    main()
