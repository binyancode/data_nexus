"""Build the Data Nexus PPTX (18 slides): PNG images + example-driven speaker notes."""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent
sys.path.insert(0, str(BASE))
import html_to_pptx as h  # noqa: E402

DOC = BASE.parent / "doc"
HTML = DOC / "nexus_platform_ppt.html"
PNG_DIR = DOC / "svg_nexus"
OUT = DOC / "Data Nexus 平台设计.pptx"

BG = RGBColor(0x14, 0x1E, 0x30)
TEXT = RGBColor(0xEA, 0xF0, 0xF6)
MUTED = RGBColor(0xB0, 0xBE, 0xC5)

slides_data = h.extract_slides(HTML)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def textbox(slide, l, t, w, ht, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, ht)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tb


for sd in slides_data:
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    if sd["is_title_slide"]:
        cover = PNG_DIR / "cover_bg.png"
        if cover.exists():
            slide.shapes.add_picture(str(cover), 0, 0,
                                     width=Inches(13.333), height=Inches(7.5))
        textbox(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.4),
                sd["title"], 54, RGBColor(0xEA, 0xF2, 0xFF), bold=True, align=PP_ALIGN.CENTER)
        if sd["subtitle"]:
            textbox(slide, Inches(1.4), Inches(4.15), Inches(10.5), Inches(1.2),
                    sd["subtitle"], 18, RGBColor(0x9F, 0xC0, 0xE8), align=PP_ALIGN.CENTER)
    else:
        textbox(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.7),
                sd["title"], 24, TEXT, bold=True)
        if sd["subtitle"]:
            textbox(slide, Inches(0.55), Inches(0.92), Inches(12.2), Inches(0.5),
                    sd["subtitle"], 13, MUTED)
        png = PNG_DIR / f"slide_{sd['index']:02d}.png"
        if png.exists():
            w = Inches(12.7)
            left = Inches((13.333 - 12.7) / 2)
            slide.shapes.add_picture(str(png), left, Inches(1.5), width=w)

    textbox(slide, Inches(11.7), Inches(7.02), Inches(1.4), Inches(0.35),
            f"{sd['index']:02d} / {len(slides_data):02d}", 10, MUTED, align=PP_ALIGN.RIGHT)


NOTES = {
    1: (
        "【第1页 · 封面】\n"
        "主标题「Data Nexus」，副标题「一张图 · 一个协议 · 一种查询 —— 联结一切数据、智能体与行动」。\n"
        "一句话：用户用大白话提问，系统自动去各种数据源（数据库、Excel/文件、知识库，甚至别的 AI 智能体）找答案、"
        "拼起来给结论，必要时还能顺手执行动作。\n"
        "全篇路线：为什么做（第2页）→ 分层架构（第3页）→ 系统蓝图（第4页）→ 逐个部件（第5-9页，含自动建本体）→ SQG（第10页）"
        "→ 执行四段：总览 + 编译/优化/协调/生成（第11-15页）→ 落地/价值/定位/深挖（第16-19页）。"
    ),
    2: (
        "【第2页 · 要解决什么问题（痛点 → 目标）】\n"
        "图上：左边 6 类数据源汇聚到中间「Data Nexus」大圆；右上橙框是用户一句话提问，右下绿框是系统自动产出的结果。\n"
        "痛点：老板问「华东上季度为什么下滑？并安排复盘」。今天要人肉打开数据库查数、翻 Excel、搜文档、找人归因、"
        "手动建复盘任务——答案散在一堆系统，每接一个新源都要单独开发，AI Agent 又是另一套。\n"
        "目标：有了 Data Nexus，用户只需一句话，系统自动：① 找到能回答的源、跨源查齐；② 综合结论并标明出处；"
        "③ 需要时顺手执行（自动派复盘任务）。\n"
        "一句话：把「人肉跨系统」变成「提一句话」。"
    ),
    3: (
        "【第3页 · 分层架构：Data Nexus 逻辑上分几层】\n"
        "这是架构思路总览：把整套系统在逻辑上分成五层，每层只做一件事、只依赖下一层的接口；后面各页就是逐层放大。\n"
        "· L5 体验层：对话 · 分析/编排画布 · 渲染（表/图/图谱/血缘/回执）——用户看得见的部分。\n"
        "· L4 意图层：Intent → SQG 编译器，把用户一句话翻成结构化查询指令。\n"
        "· L3 认知层（引擎核心）：优化器 选源规划 · 协调器 并行执行 · 生成器 合成裁决（即第11–15页那四段）。\n"
        "· L2 语义层：本体 + 能力清单（Concept · Binding · Capabilities，见第5、6页）——把业务概念映射到物理源。\n"
        "· L1 接入层：数据 Resolver · Agent Resolver · 动作 Resolver（见第7、8页）——统一抹平所有源/智能体/动作。\n"
        "右侧竖条「治理内生 · 权限/血缘/成本」：不是单独一层，而是竖切面贯穿全层——权限和血缘随每个概念一起挂在本体上"
        "（第5页 Concept 的 policy/provenance）。\n"
        "落地：按层切模块，每层只依赖下一层接口；L3 认知层是研发重心。这一页先给全貌，第4页起逐块放大。"
    ),
    4: (
        "【第4页 · 系统蓝图：所有部件如何联合运作】\n"
        "全局总图，后面每页都是放大它的一块。分三层：\n"
        "▶ 上层「知识层 · 本体」（静态）：Concept（业务名词，第5页）、Binding（概念映射到真实表，第6页）。\n"
        "▶ 中层「运行引擎」（动态）：用户提问 → 编译器（自然语言 + Concept → 查询指令 SQG）→ SQG（第10页）→ "
        "优化器 Optimizer（选源）→ 协调器 Coordinator（并行执行 + 合并 + 裁决）→ 生成器 Generator（写答案 + 血缘）→ 答案。"
        "引擎四段细节见第11–15页。\n"
        "▶ 下层「能力层 · Resolver」（静态）：数据源/Agent/动作 Resolver，详见第7、8页。\n"
        "两条虚线是「供给」：编译时向上查 Concept/Binding；执行时向下调用 Resolver。"
    ),
    5: (
        "【第5页 · Concept：业务名词长什么样】\n"
        "（放大蓝图上层的 Concept）图上：左边 concept.json 代码，右边逐字段说明。\n"
        "Concept = 业务名词：销售额=指标、地区=维度、客户=实体。只讲「是什么」，不管数据存哪。\n"
        "以「销售额」为例：· id 唯一标识。· kind 类型（实体/属性/关系/指标/维度），本例是指标。· semantics 业务含义。"
        "· type 数据类型。· bindings（★）指向哪些真实源，可多个（跨源融合的支点，见第6页）。· policy 权限。· provenance 血缘。\n"
        "要点：一切本体元素都长成同一个形状，只是 kind 不同。"
    ),
    6: (
        "【第6页 · Binding：概念怎么对接物理表】\n"
        "（放大蓝图上层的 Binding）图上：左边三层（Concept → Binding → 真实表），右边拼出来的 SQL。\n"
        "Concept 只讲业务、不含物理，靠 Binding（绑定/映射）当桥落到真实表：销售额 = SUM(fact_sales.amount)、"
        "地区 = dim_region.region_name，以及怎么 join。用户问「华东上季度销售额」，编译器靠 Binding 拼出右边那段 SQL。\n"
        "三条要点：① 换库/加源只改 Binding，Concept 不动；② 同一概念可绑多个源 = 跨源融合；③ 对接不一定是表——"
        "Agent 对接的是「提问模板」，文件对接的是路径。\n"
        "那 Concept 和 Binding 谁来建？——第9页，能自动生成。"
    ),
    7: (
        "【第7页 · Resolver：谁去执行（源 = 智能体 = 动作）】\n"
        "（放大蓝图下层）图上：最上面统一接口「interface Resolver」（三方法 capabilities() / plan() / resolve()）；"
        "下面三个框各指向接口——数据 Resolver、Agent Resolver、动作 Resolver。\n"
        "最关键的一步「抹平」：数据源、AI 智能体、执行动作，三种东西长成同一个接口。\n"
        "· 数据 Resolver：resolve() = 生成查询取数。· Agent Resolver：resolve() = 用自然语言问它，返回答案+证据。"
        "· 动作 Resolver：resolve() = 执行副作用（派任务），返回回执。\n"
        "三方法：capabilities() 报能力清单（能答哪些概念、成本、时效、信任分）；plan() 评估；resolve() 干活。\n"
        "底部核心句：优化器眼里没有「源」和「Agent」之分，只有一堆报了能力清单、竞标同一段查询的 Resolver。分类见第8页。"
    ),
    8: (
        "【第8页 · Resolver 5 类清单】\n"
        "图上：一张 4 列表格（源类型/性质/交互方式/说明），5 行：① 关系库 SQL/PG（被动·确定）② 向量库/知识库"
        "（被动·近似 RAG）③ Fabric Data Agent★（主动·黑盒，自然语言 ask()）④ REST/SaaS API（被动·外部）"
        "⑤ 动作：写回/审批/触发★（主动·有状态）。\n"
        "「性质」：被动 = 你用查询语言去取数、结果确定；主动 = 你把需求交给它、它自己去办。\n"
        "重点：主动源（第 3、5 行）是本平台差异化——把智能体和写回动作也纳入同一套 Resolver 接口。"
    ),
    9: (
        "【第9页 · Concept 和 Binding 怎么自动生成】\n"
        "第5、6页的 Concept 和 Binding 不能只靠人手工输入，那样本体永远建不起来。做法是让 Resolver 去探测源、抽样数据自动推断。\n"
        "图上：四步流程——① 探测（describe() 让源自己报出有哪些表、字段，再抽几行样例）→ ② 理解（看列名/类型/去重数/"
        "样例值/外键）→ ③ 生成候选（Concept + Binding + 置信度）→ ④ 确认（高置信自动收，低置信标红）。\n"
        "销售例子：探测到 fact_sales(amount, region_id, order_date)、dim_region(region_id, region_name) 和外键，就自动推断出："
        "销售额=指标、地区=维度；Binding 销售额=SUM(fact_sales.amount)、地区=dim_region.region_name；join 也自动填好。\n"
        "跨源合并：SQL 的「销售额」和 Agent 懂的「销售额」，系统识别是同一个，自动合并成一个概念、绑两个源 = 跨源融合。\n"
        "关键：半自动——机器生成候选，人只确认/改名剩下的；遇到含义歧义、度量可加性、敏感数据交给人（敏感数据只抽脱敏统计）。\n"
        "不同源探测方式不同：SQL 最全；文件靠列名匹配猜关联；Data Agent 是黑盒，改成直接问它「你能答哪些概念」让它自报能力清单。"
    ),
    10: (
        "【第10页 · SQG（语义查询图）：问题被翻成的统一查询指令】\n"
        "SQG = Semantic Query Graph = 语义查询图，通俗叫「查询指令」。\n"
        "图上：顶部流水线 提问→编译→SQG→规划→执行→合并；中间算子；底部结论。\n"
        "算子：SELECT 取属性；FILTER 维度约束（只看华东）；AGGREGATE 指标计算（对销售额求和）；ASK★ 把一段交给 Agent 回答；"
        "ACT★ 执行动作/写回；JOIN 跨源合并（据 relation 自动生成，同源则融进下推 SQL、跨源才在内存里做）。\n"
        "底部结论：问答=ASK、写回=ACT 都是同级算子，所以「分析」和「行动」用同一套引擎——既能答也能做。\n"
        "它具体长什么样、怎么被执行，见第11–15页。"
    ),
    11: (
        "【第11页 · 执行总览：一次提问，四段引擎跑完】\n"
        "从一个稍复杂的问题开始：「华东上季度毛利为何下滑、跟华南比差在哪？给份说明，并把复盘任务派给区域负责人。」"
        "——它既要查数、又要跨区对比、还要归因、最后派任务。\n"
        "图上：顶部是这句提问，下面四段引擎横向串起来 —— ① 编译器 → ② 优化器 → ③ 协调器 → ④ 生成器；"
        "橙色高亮的优化器是「全流程的大脑」。\n"
        "· ① 编译器：把问题拆成一张 6 节点的 DAG（n1 华东毛利、n2 华南毛利、n3 毛利趋势、n4 佐证、n5 归因、n6 派单），只说「问什么」。\n"
        "· ② 优化器：给每个节点选中一个 Resolver，把同源节点融成一条 SQL 一次下推、跨源节点留占位符，并透传当前用户——决定「怎么最省地跑」。\n"
        "· ③ 协调器：按依赖分波，能并行就并行；上游结果运行时回填下游；最后合并、冲突按信任分裁决。\n"
        "· ④ 生成器：写出华东↓ vs 华南↑ 对比结论 + 归因 + 证据 + 血缘/信任分，并附工单回执。\n"
        "这一页是总览；第12–15页把这四段各自拆到实现级。一句话：同一套流水线跑通「查数 → 对比 → 归因 → 派任务」，"
        "换库/加源只写一个 Resolver、四段引擎不改。"
    ),
    12: (
        "【第12页 · 编译器：把问题拆成语义查询图（逻辑 SQG）】\n"
        "承接第11页，放大 ① 编译器。图左：逻辑 SQG 不是一段文字，而是一张 DAG（有向无环图）——6 个节点 + 依赖边。\n"
        "· n1 华东毛利、n2 华南毛利、n3 华东毛利趋势、n4 佐证文档 四个无依赖，可先并行；· n5 归因(ASK) 依赖 n1·n2·n3·n4"
        "（要四个数字/文档到齐才能解释）；· n6 派复盘任务(ACT) 依赖 n5（说明就是任务内容）。\n"
        "图右：每个节点「问什么」—— op（算子）+ concept（业务概念）+ 约束。比如 n1 = AGGREGATE metric.gross_margin，"
        "地区=华东、期间=2024Q1。\n"
        "三条铁律：① 只描述「问什么 + 谁依赖谁」，用的是业务概念 metric.gross_margin，不是某张表某个源；"
        "② 不含任何 Resolver/SQL/源信息，这一层压根不知道毛利存在数仓还是别处；③ 依赖是唯一的执行约束——无依赖的能并行，有依赖的必须等上游。\n"
        "关键：{n} 是占位符，逻辑层只记「谁依赖谁」，真正的值要执行时才有。把这张逻辑图变成物理计划，是下一页优化器的活。"
    ),
    13: (
        "【第13页 · 优化器：逻辑 SQG → 物理执行计划（本篇重点）】\n"
        "承接第12页，这是全流程的大脑。优化器做三件事，把「问什么」的逻辑图变成「怎么最省地跑」的物理计划：\n"
        "① 选源：给每个节点让能覆盖的 Resolver 竞标，按信任/成本/时效选中一个。\n"
        "② 同源融合下推（本页最亮点）：n1 华东毛利、n2 华南毛利、n3 趋势都在数仓 → 优化器把它们合并成一条 SQL"
        "（SELECT region, month, SUM(gross_margin) … WHERE region IN(华东,华南) GROUP BY region, month），一次往返就全拿回；"
        "本来要问 3 次变成 1 次，过滤/聚合都交给数仓算。\n"
        "③ 跨源留回填：n4 知识库检索、n5 销售 Agent、n6 工单 各自独立调用；依赖别人的 n5·n6 先留 {n} 占位符，运行时由协调器回填。\n"
        "三条原则：· 能下推就下推（别把全量拉回来自己算）；· 同源就融合（华东华南趋势合成一条 SQL；若分处不同源 → 各自取回、"
        "在内存里对齐比较，即内存 JOIN）；· 跨源留回填 + 透传 as_user（凡声明支持用户级权限的 dwh/agent/ticket 都打 user_scoped，"
        "权限在源头强制；n4 公共文档除外）。\n"
        "一句话：优化器决定「选哪个源、什么下推给源、什么留到内存、谁先谁后」——同样的问题，跑得快不快全看它。"
    ),
    14: (
        "【第14页 · 协调器：按依赖分波并行 + 回填 + 合并裁决】\n"
        "承接第13页：优化器给了物理计划，协调器负责真正跑完这张 DAG。它按依赖分成几波，能并行的并行：\n"
        "· 第1波（n1·n2·n3·n4 无依赖）：一条 SQL 一次从数仓拿回 华东毛利1200万(-15%)、华南1600万(+4%)、华东趋势480/420/300；"
        "同时并行做知识库检索（命中《价格政策》《涨价通报》）。\n"
        "· 第2波（n5 归因）：把 n1–n4 的结果回填进 prompt（毛利、华南、趋势、证据），再调销售 Agent。"
        "返回：华东主因 = 大客户 Q1 减采30% + 原材料涨价；华南靠新品拉动未受冲击（附引用证据）。\n"
        "· 第3波（n6 派单）：把 n5 的说明回填进任务内容，再调工单系统。返回：复盘工单 #4521 已派华东区域负责人张三。\n"
        "执行完做两件事：· 合并：把 n1–n5 按主题（华东·毛利·上季度、华东 vs 华南）拼成一份，n6 回执附末尾。· 裁决：若某源报的毛利"
        "与 n1（数仓·信任0.98）冲突，以数仓为准，败者数字丢弃、文字保留，裁决写入血缘。\n"
        "要点：占位符 {n1}…{n5} 在各波跑完后依次变真值——这就是「分波」的意义（下游必须等上游）；n6 工单回执是行动结果、"
        "只作附加输出，不参与按信任分裁决（裁决只针对会打架的数字）。"
    ),
    15: (
        "【第15页 · 生成器：把结果写成答案 + 血缘 + 回执】\n"
        "承接第14页：协调器合并后的结果交给生成器，组织成最终答复。\n"
        "左侧「最终答复」：· 结论——华东上季度毛利1200万(-15%)、华南1600万(+4%)，华东明显掉队、华南逆势微增；"
        "· 趋势——华东逐月480→420→300；· 归因——① 大客户Q1减采30% ② 原材料涨价，华南靠新品拉动未受同等冲击；"
        "· 证据——《区域价格政策》《原材料涨价通报》两篇；· 行动已完成——复盘工单 #4521 已派张三。\n"
        "右侧「血缘·信任·权限」：每个数字挂着出处（毛利/趋势←数仓 0.98、归因←销售 Agent 0.82、证据←知识库、工单←工单系统回执）；"
        "裁决记录（毛利以数仓为准）；权限（as_user 全程透传，行级权限在源头强制——看到的就是你有权看的）。\n"
        "好处：答案可回溯（每个数字点开见源）、分析即行动（归因完顺手派单）、加新源只写一个 Resolver 这套流程不改。\n"
        "一句话：从「提一句话」到「给结论 + 顺手把复盘任务办了」，一次答完、全程留痕。"
    ),
    16: (
        "【第16页 · 落地路径：分阶段建】\n"
        "图上：一条时间轴，五个里程碑 P0→P4。\n"
        "· P0 骨架：本体 + Resolver 接口 + 查询指令编译器；验证一个 SQL 源能问答。· P1 联邦：加向量库 Resolver + 合并/裁决；"
        "验证跨源融合。· P2 智能体源★（最大亮点）：加 ASK 算子 + 接一个 Fabric Data Agent；验证「Agent = 源」。"
        "· P3 行动：加 ACT 算子 + 动作 Resolver；验证「分析 → 行动」打通。· P4 自演化：让 Resolver 探测源自动生成本体"
        "（即第9页那套）。\n"
        "建议：每阶段都要有可演示闭环；P0 就把三大接口（Concept/Resolver/查询指令）定死，之后只加实现、不改范式。\n"
        "风险点：本体自动生成的准确率、Agent 返回证据的结构化、跨源主键对齐。"
    ),
    17: (
        "【第17页 · 价值：这套设计带来什么】\n"
        "图上：五张卡片 ①~⑤ + 底部判据。\n"
        "① 加源零范式：接一个新系统 = 写一个 Resolver，引擎不改。② 一种查询：查数据/血缘/权限/问 Agent 全是同一种查询，"
        "不用维护多套 API。③ 分析即行动：ASK 和 ACT 同级，洞察到执行在同一引擎里打通。④ 治理内生：权限和血缘随概念一起管"
        "（第5页 policy/provenance），可追溯是默认。⑤ Agent 无缝：多智能体联邦是 Resolver 竞标机制的自然子集。\n"
        "底部判据：加任何新东西，本质只是「本体里加一个概念，或加一个 Resolver」，范式不变。"
    ),
    18: (
        "【第18页 · 三个战略定位】\n"
        "图上：三列方案 A/B/C + 两条横幅。A · 联邦语义中台：跨源统一问答 + 血缘治理；对标 Palantir、dbt；偏重、B 端。"
        "B · Agent 编排总控：把 Fabric、MCP、专家 Agent 编排起来；踩中 Agent 热点；轻、快。C · 自助分析 Copilot（★）："
        "业务人员自然语言跨源分析 + 自动可视化；是当前 PowerBI Copilot 的超集；落地最短。\n"
        "建议：C 起步（复用现有资产最多）→ 架构按 A 分层 → 内置 B 的「Agent 即源」做核心差异化。\n"
        "可复用资产：Orchestrator→优化器、ToolProvider→Resolver、FrontChannel、Renderer、Credential、API 权限体系。"
    ),
    19: (
        "【第19页 · 后续可深挖的四个方向】\n"
        "图上：2×2 四张卡片。1 · 本体存储选型：属性图 vs RDF vs 关系库模拟；P0 用 SQL 模拟起步。2 · 查询指令的完整算子代数 + "
        "编译器：把算子定义完整，设计 NL→查询指令的受约束生成。3 · Resolver 竞标/路由算法：用「成本-信任-时效」加权模型"
        "决定哪个源中标、冲突如何裁决。4 · 自动建本体的准确率：把结构探测 + 抽样 + LLM 推断的候选做得更准，减少人工确认。\n"
        "任选其一，可展开到接口 / 伪代码级。"
    ),
}

for i, slide in enumerate(prs.slides, start=1):
    pass  # 讲稿备注已按需求移除，不写入 notes

prs.save(str(OUT))
print("Saved:", OUT, "| slides:", len(slides_data))
